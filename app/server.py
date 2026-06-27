"""
Lokaler HTTP-Server (nur 127.0.0.1). Endpoints:

  GET  /ping                          – Health + Stats
  GET  /search?q=...&limit=...        – Suche (Pages + Files)

  POST /ingest                        – Die besuchte SharePoint-Seite (mit HTML + Links)
         Body JSON: {url, title, content, html, hostname, links:[{href,text}]}
         Antwort : {ok, id, fetch:[url,...]}   – Liste der zu crawlenden Links

  POST /ingest_page                   – Ergebnis einer gecrawlten HTML-Unterseite
         Body JSON: {url, title, content, html, hostname, source_url}

  POST /ingest_file?url=...&source_url=...&filename=...&mime_type=...
         Body: Binär-Rohdaten der Datei
"""

from __future__ import annotations

import hashlib
import json
import logging
import re
import threading
from http.server import BaseHTTPRequestHandler, ThreadingHTTPServer
from pathlib import Path
from urllib.parse import parse_qs, unquote, urlparse

from db import Database

log = logging.getLogger("sp-scraper.server")

# ---- Konfiguration -----------------------------------------------------

_CRAWLABLE_MIMES_HTML = {"text/html", "application/xhtml+xml"}
_CRAWLABLE_MIMES_FILE_PREFIXES = (
    "application/pdf",
    "application/vnd.openxmlformats-officedocument",  # docx, pptx, xlsx
    "application/vnd.ms-",                            # ppt, doc, xls (alt)
    "application/msword",
    "application/vnd.oasis.opendocument",             # odt, odp, ods
    "text/plain",
    "text/csv",
    "text/markdown",
    "application/rtf",
    "application/zip",
    "image/",
)

# Maximale Datei-Größe pro Download (Bytes)
MAX_FILE_BYTES = 100 * 1024 * 1024  # 100 MB

# Maximal wie viele Links pro Seite crawlen
MAX_LINKS_PER_PAGE = 200

# ---- Modul-Status ------------------------------------------------------

_db: Database | None = None
_paused = False
_pause_lock = threading.Lock()


def set_paused(value: bool) -> None:
    global _paused
    with _pause_lock:
        _paused = bool(value)


def is_paused() -> bool:
    with _pause_lock:
        return _paused


# ---- Helpers -----------------------------------------------------------


def _normalize_url(u: str) -> str:
    """Fragment entfernen, Whitespace strippen. Nur http(s) durchlassen."""
    u = (u or "").strip()
    if not u:
        return ""
    try:
        p = urlparse(u)
    except Exception:
        return ""
    if p.scheme not in ("http", "https"):
        return ""
    # Fragment raus
    p = p._replace(fragment="")
    return p.geturl()


def _safe_filename(name: str) -> str:
    name = (name or "").strip() or "datei"
    name = re.sub(r"[^\w.\- ]+", "_", name, flags=re.UNICODE)
    return name[:200]


def _pick_mime(mime: str) -> str | None:
    """Ist dieser MIME-Typ interessant? Gibt 'html' oder 'file' oder None."""
    if not mime:
        return None
    m = mime.split(";")[0].strip().lower()
    if m in _CRAWLABLE_MIMES_HTML:
        return "html"
    if any(m.startswith(p) for p in _CRAWLABLE_MIMES_FILE_PREFIXES):
        return "file"
    return None


# ---- HTTP Handler ------------------------------------------------------


class _Handler(BaseHTTPRequestHandler):
    # CORS: die Extension postet von beliebigen Origins
    def _send_cors(self):
        self.send_header("Access-Control-Allow-Origin", "*")
        self.send_header("Access-Control-Allow-Methods", "POST, GET, OPTIONS")
        self.send_header("Access-Control-Allow-Headers", "Content-Type")
        # Chrome Private Network Access: Requests von (sicheren) Webseiten auf
        # lokale Adressen (127.0.0.1) erfordern diese Freigabe im Preflight.
        self.send_header("Access-Control-Allow-Private-Network", "true")

    def log_message(self, fmt, *args):
        log.debug("%s - %s", self.address_string(), fmt % args)

    def _json(self, status: int, payload):
        body = json.dumps(payload).encode("utf-8")
        self.send_response(status)
        self.send_header("Content-Type", "application/json; charset=utf-8")
        self.send_header("Content-Length", str(len(body)))
        self._send_cors()
        self.end_headers()
        self.wfile.write(body)

    def _read_json(self) -> dict:
        length = int(self.headers.get("Content-Length", "0"))
        raw = self.rfile.read(length) if length > 0 else b"{}"
        return json.loads(raw.decode("utf-8") or "{}")

    def _read_bytes(self, max_bytes: int) -> bytes:
        length = int(self.headers.get("Content-Length", "0"))
        if length <= 0:
            return b""
        if length > max_bytes:
            raise ValueError(f"Datei zu groß: {length} bytes (max {max_bytes})")
        return self.rfile.read(length)

    def do_OPTIONS(self):
        self.send_response(204)
        self._send_cors()
        self.end_headers()

    # ----- GET --------------------------------------------------------
    def do_GET(self):
        assert _db is not None
        parsed = urlparse(self.path)
        if parsed.path == "/ping":
            s = _db.stats()
            s["paused"] = is_paused()
            return self._json(200, s)
        if parsed.path == "/search":
            params = parse_qs(parsed.query)
            q = (params.get("q") or [""])[0]
            limit = int((params.get("limit") or ["50"])[0])
            return self._json(200, {"results": _db.search(q, limit=limit)})
        if parsed.path == "/page":
            params = parse_qs(parsed.query)
            try:
                pid = int((params.get("id") or [""])[0])
            except ValueError:
                return self._json(400, {"error": "id erforderlich"})
            page = _db.get_page(pid)
            if page is None:
                return self._json(404, {"error": "page not found"})
            return self._json(200, page)
        if parsed.path == "/file":
            params = parse_qs(parsed.query)
            try:
                fid = int((params.get("id") or [""])[0])
            except ValueError:
                return self._json(400, {"error": "id erforderlich"})
            f = _db.get_file(fid)
            if f is None:
                return self._json(404, {"error": "file not found"})
            return self._json(200, f)
        if parsed.path == "/links":
            params = parse_qs(parsed.query)
            try:
                pid = int((params.get("page_id") or [""])[0])
            except ValueError:
                return self._json(400, {"error": "page_id erforderlich"})
            return self._json(200, {"links": _db.links_for_page(pid)})
        return self._json(404, {"error": "not found"})

    # ----- POST -------------------------------------------------------
    def do_POST(self):
        assert _db is not None
        parsed = urlparse(self.path)

        try:
            if parsed.path == "/ingest":
                return self._handle_ingest()
            if parsed.path == "/ingest_page":
                return self._handle_ingest_page()
            if parsed.path == "/ingest_file":
                return self._handle_ingest_file(parsed.query)
            if parsed.path == "/delete":
                return self._handle_delete()
            if parsed.path == "/pause":
                return self._handle_pause()
            return self._json(404, {"error": "not found"})
        except ValueError as e:
            return self._json(400, {"error": str(e)})
        except Exception as e:
            log.exception("Interner Fehler im POST %s", parsed.path)
            return self._json(500, {"error": f"internal: {e}"})

    # ----- handlers ---------------------------------------------------
    def _handle_ingest(self):
        if is_paused():
            return self._json(202, {"ok": True, "skipped": "paused", "fetch": []})
        data = self._read_json()

        url = _normalize_url(data.get("url") or "")
        title = (data.get("title") or "").strip()
        content = (data.get("content") or "").strip()
        html = data.get("html") or ""
        hostname = (data.get("hostname") or "").strip()
        links = data.get("links") or []

        if not url or not content:
            raise ValueError("url und content sind erforderlich")

        # Safety-Limits
        if len(content) > 5_000_000:
            content = content[:5_000_000]
        if len(html) > 15_000_000:
            html = html[:15_000_000]

        pid = _db.upsert_page(
            url=url, title=title, content=content, html=html,
            hostname=hostname, source="visit",
        )

        # Links normalisieren + deduplizieren
        norm_links: list[dict] = []
        seen: set[str] = set()
        for li in links:
            href = _normalize_url((li or {}).get("href") or "")
            if not href or href == url or href in seen:
                continue
            seen.add(href)
            norm_links.append({"href": href, "text": (li.get("text") or "")})
            if len(norm_links) >= MAX_LINKS_PER_PAGE:
                break

        to_crawl = _db.record_links(pid, norm_links)

        return self._json(200, {
            "ok": True,
            "id": pid,
            "fetch": to_crawl,
            "source_page_id": pid,
        })

    def _handle_ingest_page(self):
        if is_paused():
            return self._json(202, {"ok": True, "skipped": "paused"})
        data = self._read_json()

        url = _normalize_url(data.get("url") or "")
        source_url = _normalize_url(data.get("source_url") or "")
        title = (data.get("title") or "").strip()
        content = (data.get("content") or "").strip()
        html = data.get("html") or ""
        hostname = (data.get("hostname") or "").strip()

        if not url or (not content and not html):
            raise ValueError("url und (content oder html) sind erforderlich")

        if len(content) > 5_000_000:
            content = content[:5_000_000]
        if len(html) > 15_000_000:
            html = html[:15_000_000]

        pid = _db.upsert_page(
            url=url, title=title, content=content, html=html,
            hostname=hostname, source="crawl",
        )

        # Quelle markieren
        if source_url:
            src = self._source_page_id(source_url)
            if src is not None:
                _db.mark_link(src, url, status="fetched_page", target_page_id=pid)

        return self._json(200, {"ok": True, "id": pid})

    def _handle_ingest_file(self, query: str):
        if is_paused():
            return self._json(202, {"ok": True, "skipped": "paused"})
        params = parse_qs(query)
        url = _normalize_url((params.get("url") or [""])[0])
        source_url = _normalize_url((params.get("source_url") or [""])[0])
        filename = _safe_filename(unquote((params.get("filename") or [""])[0]))
        mime_type = (params.get("mime_type") or [""])[0]

        if not url:
            raise ValueError("url ist erforderlich")

        body = self._read_bytes(MAX_FILE_BYTES)
        if not body:
            raise ValueError("leerer Body")

        sha = hashlib.sha256(body).hexdigest()
        # Erweiterung aus filename ableiten (oder mime)
        ext = Path(filename).suffix.lower().lstrip(".") or _mime_to_ext(mime_type)
        ext = re.sub(r"[^a-z0-9]+", "", ext)[:8] or "bin"
        subdir = _db.files_dir / sha[:2]
        subdir.mkdir(parents=True, exist_ok=True)
        path = subdir / f"{sha}.{ext}"
        if not path.exists():
            path.write_bytes(body)

        extracted = _extract_text_if_possible(path, mime_type)

        fid = _db.upsert_file(
            url=url,
            sha256=sha,
            filename=filename or path.name,
            mime_type=mime_type,
            size=len(body),
            local_path=str(path),
            extracted_text=extracted,
        )

        if source_url:
            src = self._source_page_id(source_url)
            if src is not None:
                _db.mark_link(src, url, status="fetched_file", target_file_id=fid)

        return self._json(200, {"ok": True, "id": fid, "local_path": str(path)})

    def _handle_delete(self):
        data = self._read_json()
        kind = (data.get("kind") or "").strip()
        try:
            item_id = int(data.get("id"))
        except (TypeError, ValueError):
            raise ValueError("id erforderlich")
        if kind == "page":
            _db.delete_page(item_id)
        elif kind == "file":
            _db.delete_file(item_id)
        else:
            raise ValueError("kind muss 'page' oder 'file' sein")
        return self._json(200, {"ok": True})

    def _handle_pause(self):
        data = self._read_json()
        value = bool(data.get("value"))
        set_paused(value)
        return self._json(200, {"ok": True, "paused": is_paused()})

    # ----- util -------------------------------------------------------
    def _source_page_id(self, source_url: str) -> int | None:
        assert _db is not None
        return _db.page_id_by_url(source_url)


# --------------------------------------------------------------------- text


def _mime_to_ext(mime: str) -> str:
    m = (mime or "").split(";")[0].strip().lower()
    return {
        "application/pdf": "pdf",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
        "application/vnd.ms-powerpoint": "ppt",
        "application/msword": "doc",
        "application/vnd.ms-excel": "xls",
        "text/plain": "txt",
        "text/csv": "csv",
        "text/markdown": "md",
        "application/rtf": "rtf",
        "application/zip": "zip",
    }.get(m, "")


def _extract_text_if_possible(path: Path, mime: str) -> str:
    """
    Best-Effort-Textextraktion für die FTS-Suche. Nutzt nur stdlib + optionale
    Bibliotheken, falls installiert (pypdf, python-docx, python-pptx).
    Bei Fehlern: leerer String.
    """
    try:
        m = (mime or "").split(";")[0].strip().lower()
        if m.startswith("text/") or path.suffix.lower() in {".txt", ".csv", ".md"}:
            return path.read_text(encoding="utf-8", errors="replace")[:2_000_000]

        if m == "application/pdf" or path.suffix.lower() == ".pdf":
            try:
                import pypdf  # type: ignore
                reader = pypdf.PdfReader(str(path))
                parts = []
                for page in reader.pages[:500]:
                    try:
                        parts.append(page.extract_text() or "")
                    except Exception:
                        pass
                return "\n".join(parts)[:2_000_000]
            except Exception:
                return ""

        if path.suffix.lower() == ".docx":
            try:
                import docx  # python-docx
                d = docx.Document(str(path))
                return "\n".join(p.text for p in d.paragraphs)[:2_000_000]
            except Exception:
                return ""

        if path.suffix.lower() == ".pptx":
            try:
                from pptx import Presentation
                p = Presentation(str(path))
                parts = []
                for slide in p.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            parts.append(shape.text)
                return "\n".join(parts)[:2_000_000]
            except Exception:
                return ""

        if path.suffix.lower() == ".xlsx":
            try:
                import openpyxl  # noqa
                wb = openpyxl.load_workbook(str(path), data_only=True, read_only=True)
                parts = []
                for ws in wb.worksheets:
                    for row in ws.iter_rows(values_only=True):
                        parts.append(" ".join("" if c is None else str(c) for c in row))
                return "\n".join(parts)[:2_000_000]
            except Exception:
                return ""

    except Exception:
        return ""
    return ""


# --------------------------------------------------------------------- server


def start_server_thread(db: Database, host: str = "127.0.0.1", port: int = 8765):
    global _db
    _db = db
    server = ThreadingHTTPServer((host, port), _Handler)
    thread = threading.Thread(target=server.serve_forever, name="sp-scraper-http", daemon=True)
    thread.start()
    log.info("Lokaler Server lauscht auf http://%s:%s", host, port)
    return server, thread


def run_server(db: Database, host: str = "127.0.0.1", port: int = 8765):
    """Blockierender Start (für Tests)."""
    global _db
    _db = db
    server = ThreadingHTTPServer((host, port), _Handler)
    log.info("Lokaler Server lauscht auf http://%s:%s", host, port)
    server.serve_forever()
    return server
